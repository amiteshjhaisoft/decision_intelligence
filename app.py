# Author: Amitesh Jha | iSoft | 2025-10-12
# Streamlit RAG chat — Azure Blob → Weaviate (hybrid) → (optional) Cross-Encoder rerank → Claude (grounded-only)
# Now with DI Agent: Plan → Retrieve → Validate → Act (KB-only skills) → Respond (+ exports)

from __future__ import annotations

import io, os, json, socket, ssl, tempfile, re, html, math, statistics, datetime as dt
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import urlparse

import streamlit as st

# ---- Avoid odd TLS EOFs through certain proxies/edges (must be set BEFORE importing weaviate)
os.environ.setdefault("HTTPX_DISABLE_HTTP2", "1")

# ---------------------- Optional heavy extractors (graceful if missing) ----------------------
# Images (OCR)
try:
    from PIL import Image, ExifTags  # pip install pillow
except Exception:
    Image, ExifTags = None, None
try:
    import pytesseract  # pip install pytesseract (needs tesseract binary installed)
except Exception:
    pytesseract = None

# Audio / Video (ASR)
try:
    import whisper  # pip install -U openai-whisper
except Exception:
    whisper = None
try:
    from moviepy.editor import AudioFileClip  # pip install moviepy
except Exception:
    AudioFileClip = None

# HTML parsing
try:
    from bs4 import BeautifulSoup  # pip install beautifulsoup4
except Exception:
    BeautifulSoup = None

# Office / PDF
try:
    from pypdf import PdfReader  # pip install pypdf
except Exception:
    PdfReader = None
try:
    import docx  # pip install python-docx
except Exception:
    docx = None
try:
    from pptx import Presentation  # pip install python-pptx
except Exception:
    Presentation = None
try:
    import pandas as pd  # pip install pandas openpyxl
except Exception:
    pd = None

# Text splitting
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter  # pip install langchain-text-splitters
except Exception as e:
    raise RuntimeError("Missing dependency: langchain-text-splitters") from e

# Azure + Weaviate + Embeddings + LLM
try:
    from azure.storage.blob import BlobServiceClient
except Exception as e:
    raise RuntimeError("Missing dependency: azure-storage-blob") from e

try:
    import weaviate  # v4 preferred (>=4.6.0)
except Exception as e:
    raise RuntimeError("Missing dependency: weaviate-client") from e

# Detect v4 helpers
try:
    from weaviate.classes.init import Auth
    from weaviate.classes.config import Property, DataType, Configure
    V4 = True
except Exception:
    V4 = False

try:
    from weaviate.classes.tenants import Tenant  # v4 only
except Exception:
    Tenant = None
try:
    from weaviate.classes.query import MetadataQuery
except Exception:
    MetadataQuery = None

# Embeddings (dense) + Cross-Encoder rerank (optional)
try:
    from sentence_transformers import SentenceTransformer
except Exception as e:
    raise RuntimeError("Missing dependency: sentence-transformers") from e

try:
    from sentence_transformers import CrossEncoder  # optional
except Exception:
    CrossEncoder = None

# Anthropic (Claude)
try:
    from anthropic import Anthropic
except Exception as e:
    raise RuntimeError("Missing dependency: anthropic") from e


# ---------------------- App constants ----------------------
CLASS_NAME       = "KBChunk"
CHUNK_SIZE       = 900
CHUNK_OVERLAP    = 150

# Retrieval configuration (non-LLM only)
CANDIDATES_K     = 60     # hybrid recall pool
FINAL_K_DEFAULT  = 10     # passages sent to LLM

# Anthropic model (adjust to your account/allowlist)
CLAUDE_MODEL     = "claude-3-7-sonnet-20250219"

# Bounds to avoid exploding memory on huge binaries
MAX_BYTES_IMAGE = 25 * 1024 * 1024   # 25 MB
MAX_BYTES_AUDIO = 80 * 1024 * 1024   # 80 MB
MAX_BYTES_VIDEO = 150 * 1024 * 1024  # 150 MB
MAX_BYTES_TEXT  = 50 * 1024 * 1024   # 50 MB

# File type buckets
TEXT_EXTS  = {".txt", ".md", ".log"}
CSV_EXT    = {".csv"}
JSON_EXT   = {".json"}
YAML_EXT   = {".yaml", ".yml"}
HTML_EXT   = {".html", ".htm"}
DOC_EXT    = {".docx"}
PPT_EXT    = {".pptx"}
XLS_EXT    = {".xlsx"}
PDF_EXT    = {".pdf"}
CODE_EXTS  = {
    ".py", ".sql", ".js", ".ts", ".java", ".c", ".cpp", ".cs",
    ".go", ".rb", ".rs", ".php", ".xml", ".ini", ".cfg", ".toml"
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp", ".gif"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac"}
VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".webm"}


# ---------------------- Small helpers ----------------------
def _hash(s: str) -> str:
    import hashlib as _h
    return _h.md5(s.encode("utf-8")).hexdigest()

def _truncate(s: str, limit: int = 2000) -> str:
    s = s or ""
    return s if len(s) <= limit else s[:limit] + " …"

def _safe_decode(b: bytes, encoding="utf-8") -> str:
    try:
        return b.decode(encoding, errors="ignore")
    except Exception:
        return ""

def esc(s: str) -> str:
    """Escape any user/LLM text before putting inside our HTML bubble container."""
    return html.escape(s or "")

def has_citation(text: str) -> bool:
    return bool(re.search(r"\[\d+\]", text or ""))


# --------- Parsers (each returns a TEXT representation suitable for RAG) ----------
def read_pdf_bytes(b: bytes) -> str:
    if not PdfReader:
        return ""
    try:
        reader = PdfReader(io.BytesIO(b))
        return "\n".join((page.extract_text() or "") for page in reader.pages).strip()
    except Exception:
        return ""

def read_docx_bytes(b: bytes) -> str:
    if not docx:
        return ""
    try:
        d = docx.Document(io.BytesIO(b))
        return "\n".join(p.text for p in d.paragraphs).strip()
    except Exception:
        return ""

def read_pptx_bytes(b: bytes) -> str:
    if not Presentation:
        return ""
    try:
        prs = Presentation(io.BytesIO(b))
        texts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        slide_text.append(t)
            if slide_text:
                texts.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
        return "\n\n".join(texts).strip()
    except Exception:
        return ""

def read_xlsx_bytes(b: bytes, limit_rows: int = 5000) -> str:
    if not pd:
        return ""
    try:
        buf = io.BytesIO(b)
        xl = pd.ExcelFile(buf)
        out = []
        for name in xl.sheet_names:
            df = xl.parse(name)
            if len(df) > limit_rows:
                df = df.head(limit_rows)
            out.append(f"[Sheet: {name}]\n" + df.to_csv(index=False))
        return "\n\n".join(out).strip()
    except Exception:
        return ""

def read_csv_bytes(b: bytes, limit_rows: int = 20000) -> str:
    if not pd:
        return _safe_decode(b)
    try:
        buf = io.BytesIO(b)
        df = pd.read_csv(buf, nrows=limit_rows)
        return df.to_csv(index=False)
    except Exception:
        return _safe_decode(b)

def read_html_bytes(b: bytes) -> str:
    if not BeautifulSoup:
        return _safe_decode(b)
    try:
        soup = BeautifulSoup(_safe_decode(b), "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n")
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()
    except Exception:
        return _safe_decode(b)

def read_yaml_bytes(b: bytes) -> str:
    return _safe_decode(b).strip()

def read_text_like_bytes(b: bytes, encoding: str = "utf-8") -> str:
    return _safe_decode(b, encoding)

def read_code_bytes(b: bytes) -> str:
    return _safe_decode(b)

def read_image_bytes_ocr(b: bytes, name: str) -> str:
    if not Image:
        return f"[IMAGE: {name}] (Pillow not installed; OCR unavailable)"
    try:
        with Image.open(io.BytesIO(b)) as im:
            exif_txt = ""
            try:
                exif = getattr(im, "_getexif", lambda: None)()
                if exif:
                    inv = {ExifTags.TAGS.get(k, str(k)): v for k, v in exif.items()}
                    keep = {k: inv[k] for k in list(inv)[:20]}
                    exif_txt = json.dumps(keep, default=str)
            except Exception:
                pass
            ocr_txt = ""
            if pytesseract:
                try:
                    ocr_txt = pytesseract.image_to_string(im) or ""
                except Exception as e:
                    ocr_txt = f"(OCR failed: {e})"
            else:
                ocr_txt = "(OCR skipped: pytesseract not installed)"
            return f"[IMAGE: {name}]\nEXIF: {exif_txt}\nTEXT: {ocr_txt}".strip()
    except Exception as e:
        return f"[IMAGE: {name}] (failed to open: {e})"

def transcribe_audio_bytes(b: bytes, name: str, tmp_dir: Path) -> str:
    if not whisper:
        return f"[AUDIO: {name}] (ASR skipped: whisper not installed)"
    try:
        p = tmp_dir / (Path(name).stem + ".wav")
        p.write_bytes(b)
        model = whisper.load_model("base")
        res = model.transcribe(str(p))
        txt = (res.get("text") or "").strip()
        return f"[AUDIO: {name}]\nTRANSCRIPT: {txt}" if txt else f"[AUDIO: {name}] (no speech detected)"
    except Exception as e:
        return f"[AUDIO: {name}] (ASR failed: {e})"

def transcribe_video_bytes(b: bytes, name: str, tmp_dir: Path) -> str:
    if not whisper or not AudioFileClip:
        return f"[VIDEO: {name}] (ASR skipped: whisper/moviepy not installed)"
    try:
        vid_path = tmp_dir / (Path(name).name.replace("/", "_"))
        aud_path = tmp_dir / (Path(name).stem + ".wav")
        vid_path.write_bytes(b)
        clip = AudioFileClip(str(vid_path))
        clip.audio.write_audiofile(str(aud_path), verbose=False, logger=None)
        clip.close()
        model = whisper.load_model("base")
        res = model.transcribe(str(aud_path))
        txt = (res.get("text") or "").strip()
        return f"[VIDEO: {name}]\nTRANSCRIPT: {txt}" if txt else f"[VIDEO: {name}] (no speech detected)"
    except Exception as e:
        return f"[VIDEO: {name}] (ASR failed: {e})"

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=overlap,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    return splitter.split_text(text or "")


# ---------------------- Azure loader ----------------------
def list_and_download_kb(container: str, prefix: str, connection_string: str, cache_dir: Path) -> List[Tuple[str, Path, int]]:
    """
    Download blobs to cache_dir -> [(blob_name, local_path, size_bytes)]
    Uses ETag cache to avoid re-downloading unchanged blobs.
    """
    cache_dir.mkdir(parents=True, exist_ok=True)
    bsc = BlobServiceClient.from_connection_string(connection_string)
    cont = bsc.get_container_client(container)

    results: List[Tuple[str, Path, int]] = []
    blobs = cont.list_blobs(name_starts_with=prefix or "")
    for blob in blobs:
        size = getattr(blob, "size", None)
        if not size or size <= 0:
            continue
        blob_name = blob.name
        local_path = cache_dir / blob_name.replace("/", "__")
        etag_path = cache_dir / (local_path.name + ".etag")

        etag = getattr(blob, "etag", None)
        if local_path.exists() and etag_path.exists() and etag:
            try:
                if etag_path.read_text().strip() == etag:
                    results.append((blob_name, local_path, int(size)))
                    continue
            except Exception:
                pass

        try:
            data = cont.download_blob(blob_name).readall()
            local_path.parent.mkdir(parents=True, exist_ok=True)
            local_path.write_bytes(data)
            if etag:
                etag_path.write_text(etag)
            results.append((blob_name, local_path, int(size)))
        except Exception as e:
            st.warning(f"Failed to download {blob_name}: {e}")

    return results


# ---------------------- Guess & Read (broad types) ----------------------
def guess_and_read_path(name: str, path: Path, size_bytes: int, tmp_dir: Path) -> str:
    ext = Path(name).suffix.lower()

    # Size guards by category (avoid OOM)
    if ext in IMAGE_EXTS and size_bytes > MAX_BYTES_IMAGE:
        return f"[IMAGE: {name}] (skipped: {size_bytes} bytes > {MAX_BYTES_IMAGE})"
    if ext in AUDIO_EXTS and size_bytes > MAX_BYTES_AUDIO:
        return f"[AUDIO: {name}] (skipped: {size_bytes} bytes > {MAX_BYTES_AUDIO})"
    if ext in VIDEO_EXTS and size_bytes > MAX_BYTES_VIDEO:
        return f"[VIDEO: {name}] (skipped: {size_bytes} bytes > {MAX_BYTES_VIDEO})"
    if (ext in (TEXT_EXTS | CSV_EXT | JSON_EXT | YAML_EXT | HTML_EXT | DOC_EXT | PPT_EXT | XLS_EXT | PDF_EXT | CODE_EXTS)) and size_bytes > MAX_BYTES_TEXT:
        return f"[TEXT-LIKE: {name}] (skipped: {size_bytes} bytes > {MAX_BYTES_TEXT})"

    b = path.read_bytes()

    # Text-like
    if ext in TEXT_EXTS:
        return read_text_like_bytes(b)
    if ext in CSV_EXT:
        return read_csv_bytes(b)
    if ext in JSON_EXT:
        try:
            j = json.loads(_safe_decode(b))
            return json.dumps(j, indent=2, ensure_ascii=False)
        except Exception:
            return _safe_decode(b)
    if ext in YAML_EXT:
        return read_yaml_bytes(b)
    if ext in HTML_EXT:
        return read_html_bytes(b)
    if ext in DOC_EXT:
        return read_docx_bytes(b)
    if ext in PPT_EXT:
        return read_pptx_bytes(b)
    if ext in XLS_EXT:
        return read_xlsx_bytes(b)
    if ext in PDF_EXT:
        return read_pdf_bytes(b)
    if ext in CODE_EXTS:
        return read_code_bytes(b)

    # Binary media
    if ext in IMAGE_EXTS:
        return read_image_bytes_ocr(b, name)
    if ext in AUDIO_EXTS:
        return transcribe_audio_bytes(b, name, tmp_dir)
    if ext in VIDEO_EXTS:
        return transcribe_video_bytes(b, name, tmp_dir)

    # Fallback
    return _safe_decode(b)


# ---------------------- Connectivity diagnostics ----------------------
def run_weaviate_diagnostics(base_url: str, api_key: Optional[str]) -> Dict[str, Any]:
    info: Dict[str, Any] = {"url": base_url}
    try:
        u = urlparse(base_url)
        host = u.hostname or ""
        port = u.port or (443 if u.scheme == "https" else 80)
        info["parsed"] = {"scheme": u.scheme, "host": host, "port": port}

        # DNS
        try:
            addrs = socket.getaddrinfo(host, port, proto=socket.IPPROTO_TCP)
            info["dns"] = [f"{a[4][0]}:{a[4][1]}" for a in addrs]
        except Exception as e:
            info["dns_error"] = repr(e)

        # TCP
        try:
            s = socket.create_connection((host, port), timeout=6)
            s.close()
            info["tcp_ok"] = True
        except Exception as e:
            info["tcp_error"] = repr(e)

        # HTTPS /.well-known/ready (no auth)
        try:
            import urllib.request
            ready_url = base_url.rstrip("/") + "/v1/.well-known/ready"
            ctx = ssl.create_default_context()
            with urllib.request.urlopen(ready_url, context=ctx, timeout=8) as r:
                info["ready_no_auth"] = f"HTTP {r.status}"
        except Exception as e:
            info["ready_no_auth_error"] = repr(e)

        # HTTPS /.well-known/ready (with Bearer)
        try:
            import urllib.request
            ready_url = base_url.rstrip("/") + "/v1/.well-known/ready"
            ctx = ssl.create_default_context()
            req = urllib.request.Request(ready_url, headers={"Authorization": f"Bearer {api_key}"} if api_key else {})
            with urllib.request.urlopen(req, context=ctx, timeout=8) as r:
                info["ready_with_auth"] = f"HTTP {r.status}"
        except Exception as e:
            info["ready_with_auth_error"] = repr(e)
    except Exception as e:
        info["diagnostics_error"] = repr(e)
    return info


# ---------------------- Weaviate client (v4 with safe fallbacks) ----------------------
def make_weaviate_client(url: str, api_key: Optional[str]) -> Any:
    ver = getattr(weaviate, "__version__", "unknown")
    is_v3_pkg = ver.startswith("3.")
    st.caption(f"📦 weaviate-client version: {ver} (V4 symbols: {bool(V4)})")

    def _parse(u: str):
        if not u or "://" not in u:
            raise ValueError(
                f"Invalid Weaviate URL '{u}'. Use full https URL (e.g., https://<cluster-id>.weaviate.cloud)."
            )
        p = urlparse(u)
        if not p.hostname:
            raise ValueError(f"Could not parse hostname from '{u}'")
        https = (p.scheme == "https")
        http_port = p.port or (443 if https else 80)
        return p.hostname, https, http_port

    if V4:
        auth = Auth.api_key(api_key) if api_key else None

        # 1) WCS (.weaviate.network / .weaviate.cloud)
        try:
            if (".weaviate.network" in url) or (".weaviate.cloud" in url) or ("wcs" in url):
                st.caption("🔌 connect_to_weaviate_cloud(cluster_url=...)")
                client = weaviate.connect_to_weaviate_cloud(
                    cluster_url=url,
                    auth_credentials=auth,
                    skip_init_checks=True,
                )
                client.is_connected()
                st.caption("✅ WCS connected")
                return client
        except Exception as e:
            st.warning(f"WCS connect failed (will try custom): {e}")

        # 2) Custom (include grpc args for newer v4)
        try:
            host, https, http_port = _parse(url)
            grpc_host = host
            grpc_secure = https
            grpc_port = 50051
            st.caption(
                f"🔌 connect_to_custom(http={host}:{http_port} https={https}, grpc={grpc_host}:{grpc_port} tls={grpc_secure})"
            )
            client = weaviate.connect_to_custom(
                http_host=host,
                http_port=http_port,
                http_secure=https,
                grpc_host=grpc_host,
                grpc_port=grpc_port,
                grpc_secure=grpc_secure,
                auth_credentials=auth,
                skip_init_checks=True,
            )
            client.is_connected()
            st.caption("✅ custom connected")
            return client
        except Exception as e:
            st.warning(f"Custom connect failed (v4): {e}")

    # 3) v3 fallback
    if is_v3_pkg:
        try:
            st.caption("🔌 v3 Client(url=...)")
            from weaviate import Client
            from weaviate.auth import AuthApiKey
            client = Client(url=url, auth_client_secret=AuthApiKey(api_key) if api_key else None)
            client.schema.get()
            st.caption("✅ v3 connected")
            return client
        except Exception as e:
            st.warning(f"v3 fallback failed: {e}")

    raise RuntimeError("All connection strategies failed. Check URL (must include https://), API key, and IP allow-list.")


def ensure_collection(client: Any, class_name: str = CLASS_NAME, tenancy: Optional[str] = None) -> None:
    """
    Ensure a KBChunk collection/class exists. If a tenant name is provided, ensure the tenant exists.
    """
    if V4:
        try:
            existing_names = client.collections.list_all()
            existing_names = [getattr(c, "name", c) for c in existing_names]
        except Exception:
            existing_names = []
        if class_name not in existing_names:
            mt_cfg = Configure.multi_tenancy(enabled=bool(tenancy)) if tenancy else None
            client.collections.create(
                name=class_name,
                vectorizer_config=Configure.Vectorizer.none(),
                properties=[
                    Property(name="doc_id", data_type=DataType.TEXT),
                    Property(name="source", data_type=DataType.TEXT),
                    Property(name="chunk_index", data_type=DataType.INT),
                    Property(name="text", data_type=DataType.TEXT),
                ],
                multi_tenancy_config=mt_cfg,
            )
        if tenancy:
            col = client.collections.get(class_name)
            try:
                existing_tenants = [t.name for t in col.tenants.list_all()]
            except Exception:
                existing_tenants = []
            if tenancy not in existing_tenants and Tenant is not None:
                col.tenants.create(Tenant(name=tenancy))
    else:
        schema = client.schema.get()
        existing = [c["class"] for c in schema.get("classes", [])]
        if class_name in existing:
            return
        new_class = {
            "class": class_name,
            "vectorizer": "none",
            "properties": [
                {"name": "doc_id", "dataType": ["text"]},
                {"name": "source", "dataType": ["text"]},
                {"name": "chunk_index", "dataType": ["int"]},
                {"name": "text", "dataType": ["text"]},
            ],
        }
        client.schema.create_class(new_class)


def upsert_chunks(client: Any, class_name: str, items: List[Dict[str, Any]], vectors: List[List[float]], tenancy: Optional[str] = None) -> int:
    if not items:
        return 0
    if V4:
        col = client.collections.get(class_name, tenant=tenancy) if tenancy else client.collections.get(class_name)
        if tenancy:
            # Insert one-by-one for widest compatibility with MT across v4 minors
            count = 0
            for it, vec in zip(items, vectors):
                col.data.insert(properties=it["properties"], uuid=it["id"], vector=vec)
                count += 1
            return count
        count = 0
        with col.batch.dynamic() as batch:
            for it, vec in zip(items, vectors):
                batch.add_object(properties=it["properties"], uuid=it["id"], vector=vec)
                count += 1
        return count
    else:
        batch = client.batch
        batch.configure(batch_size=64, num_workers=2)
        with batch as b:
            for it, vec in zip(items, vectors):
                b.add_data_object(
                    data_object=it["properties"],
                    class_name=class_name,
                    uuid=it["id"],
                    vector=vec,
                )
        return len(items)


# ---------------------- Hybrid retrieval → optional rerank → diversify + neighbors ----------------------
@st.cache_resource(show_spinner=False)
def get_embedder():
    return SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

@st.cache_resource(show_spinner=False)
def get_cross_encoder():
    if CrossEncoder is None:
        return None
    try:
        return CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")
    except Exception:
        return None

def _fetch_neighbors(col, class_name: str, hit: Dict[str, Any], window: int) -> List[Dict[str, Any]]:
    """Fetch ±window neighboring chunks from the same doc_id using a cheap BM25 pass."""
    if window <= 0:
        return [hit]
    doc_id = hit.get("doc_id")
    idx = hit.get("chunk_index") or 0
    want = set(range(max(0, idx - window), idx + window + 1))
    out = [hit]
    try:
        res = col.query.bm25(
            query=str(doc_id),
            limit=max(10, 3 * window + 1),
            return_properties=["doc_id", "source", "chunk_index", "text"],
        )
        for o in (getattr(res, "objects", None) or []):
            p = o.properties or {}
            if p.get("doc_id") != doc_id:
                continue
            ci = p.get("chunk_index")
            if isinstance(ci, int) and ci in want and ci != idx:
                out.append({
                    "doc_id": p.get("doc_id"),
                    "source": p.get("source"),
                    "chunk_index": ci,
                    "text": (p.get("text") or "").strip(),
                    "score": None,
                })
    except Exception:
        pass
    out = sorted({(r["doc_id"], r["chunk_index"]): r for r in out}.values(), key=lambda r: r.get("chunk_index") or 0)
    return out

def retrieve_hybrid_then_rerank(
    col,                                 # v4 collection handle (tenant already bound)
    query_text: str,
    query_vec: List[float],
    k_candidates: int = CANDIDATES_K,
    k_final: int = FINAL_K_DEFAULT,
    use_reranker: bool = True,
    max_chars_per_passage: int = 1100,
    per_source_cap: int = 3,
    neighbors_window: int = 1,
    downweight_csv: bool = True,
) -> List[Dict[str, Any]]:

    # 1) HYBRID (falls back to dense-only if hybrid unavailable)
    try:
        meta = MetadataQuery(distance=True, score=True) if MetadataQuery else None
        res = col.query.hybrid(
            query=query_text,
            vector=query_vec,
            alpha=0.5,
            limit=k_candidates,
            return_properties=["doc_id","source","chunk_index","text"],
            return_metadata=meta,
        )
    except Exception:
        res = col.query.near_vector(
            vector=query_vec,
            limit=k_candidates,
            return_properties=["doc_id","source","chunk_index","text"],
            return_metadata=MetadataQuery(distance=True) if MetadataQuery else None,
        )

    objs = getattr(res, "objects", None) or []
    rows = []
    for o in objs:
        p = o.properties or {}
        m = getattr(o, "metadata", None)
        rows.append({
            "doc_id": p.get("doc_id"),
            "source": p.get("source"),
            "chunk_index": p.get("chunk_index"),
            "text": (p.get("text") or "").strip(),
            "dist": getattr(m, "distance", None),
            "kw": getattr(m, "score", None),
        })
    if not rows:
        return []

    # 2) Lightweight blend score (lower is better)
    def is_csv_src(s: str) -> bool:
        s = (s or "").lower()
        return s.endswith(".csv") or ".csv." in s or s.endswith(".bin")
    def blend(r):
        d = r["dist"] if isinstance(r["dist"], (float, int)) else 1.0
        k = r["kw"]   if isinstance(r["kw"], (float, int)) else 0.0
        val = 0.7 * d + 0.3 * (1.0 - k)
        if downweight_csv and is_csv_src(r.get("source", "")):
            val += 0.05
        return val
    rows.sort(key=blend)

    # Optional cross-encoder rerank on the top N
    if use_reranker:
        ce = get_cross_encoder()
        if ce is not None:
            pool = rows[: max(20, k_final * 3)]
            pairs = [(query_text, r["text"]) for r in pool]
            try:
                scores = ce.predict(pairs).tolist()
                for r, s in zip(pool, scores):
                    r["ce"] = float(s)
                pool.sort(key=lambda x: x.get("ce", -1.0), reverse=True)
                rows = pool + rows[len(pool):]
            except Exception:
                pass

    # 3) Per-source cap + diversify
    per_src_counts, picked = {}, []
    for r in rows:
        src = r.get("source")
        if per_src_counts.get(src, 0) >= per_source_cap:
            continue
        picked.append(r)
        per_src_counts[src] = per_src_counts.get(src, 0) + 1
        if len(picked) >= k_final:
            break

    # 4) Neighbor stitching
    stitched = []
    for r in picked:
        stitched.extend(_fetch_neighbors(col, CLASS_NAME, r, neighbors_window))

    # 5) Pack & truncate
    out = []
    seen_keys = set()
    for r in stitched:
        key = (r.get("source"), r.get("chunk_index"))
        if key in seen_keys:
            continue
        seen_keys.add(key)
        txt = (r.get("text") or "")[:max_chars_per_passage]
        out.append({
            "doc_id": r.get("doc_id"),
            "source": r.get("source"),
            "chunk_index": r.get("chunk_index"),
            "text": txt,
            "score": r.get("dist"),
        })
    # keep at most k_final * (1 + 2*neighbors_window)
    return out[: k_final * (1 + 2*neighbors_window)]


def embed_texts(texts: List[str]) -> List[List[float]]:
    if not texts:
        return []
    model = get_embedder()
    return model.encode(texts, show_progress_bar=False, normalize_embeddings=True).tolist()


# ---------------------- Ingestion ----------------------
def ingest_from_azure_to_weaviate(connection_string: str, container: str, prefix: str,
                                  w_client: Any, tenancy: Optional[str] = None,
                                  status_cb=lambda s: None) -> int:
    ensure_collection(w_client, CLASS_NAME, tenancy=tenancy)
    cache_dir = Path(tempfile.gettempdir()) / "kb_cache_streamlit"
    tmp_dir = Path(tempfile.gettempdir()) / "kb_tmp_streamlit"
    tmp_dir.mkdir(parents=True, exist_ok=True)

    downloaded = list_and_download_kb(container, prefix, connection_string, cache_dir)
    total = 0
    for blob_name, local_path, size_bytes in downloaded:
        try:
            text = guess_and_read_path(blob_name, local_path, size_bytes, tmp_dir).strip()
            if not text:
                status_cb(f"⚠️ No text extracted from {blob_name}; skipping.")
                continue

            # CHUNK: prefix with filename so retrieval sees it as a signal
            chunks_raw = chunk_text(text)
            chunks = [f"FILE: {blob_name}\n{ch}" for ch in chunks_raw]
            if not chunks:
                continue

            doc_id = _hash(blob_name)
            items = [{
                "id": _hash(f"{doc_id}:{i}"),
                "properties": {
                    "doc_id": doc_id,
                    "source": blob_name,
                    "chunk_index": i,
                    "text": ch,
                },
            } for i, ch in enumerate(chunks)]
            vectors = embed_texts([it["properties"]["text"] for it in items])
            upserted = upsert_chunks(w_client, CLASS_NAME, items, vectors, tenancy)
            total += upserted
            status_cb(f"Indexed: {blob_name} ({size_bytes} bytes) → {upserted} chunks")
        except Exception as e:
            status_cb(f"⚠️ Failed on {blob_name}: {e}")
    return total


# ---------------------- Claude (answer-only; no retrieval use) ----------------------
def call_claude(api_key: str, system_prompt: str, user_prompt: str, stream: bool = True):
    client = Anthropic(api_key=api_key)
    if stream:
        return client.messages.stream(
            model=CLAUDE_MODEL,
            max_tokens=800,
            temperature=0.2,
            system=system_prompt,
            messages=[{"role": "user", "content": user_prompt}],
        )
    else:
        resp = client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=800,
            temperature=0.2,
            system=system_prompt,
            messages=[{"role": "user", "content": user_prompt}],
        )
        try:
            return "".join(getattr(b, "text", "") for b in (resp.content or []))
        except Exception:
            return str(resp)


# ---------------------- Prompts (modes + strictness) ----------------------
def build_system_prompt(strictness: float, mode: str) -> str:
    baseline = (
        "You are a strict, grounded assistant for a Decision Intelligence system.\n"
        "RULES:\n"
        "1) Use ONLY the provided passages (KB) for facts.\n"
        "2) Add citation markers like [1], [2] next to each claim you use.\n"
        "3) If the passages lack the answer, say: \"I don’t have enough information in the context.\" and stop.\n"
    )
    style = (
        f"Strictness: {strictness:.1f} (1.0 = quote snippets verbatim; 0.0 = freer paraphrase).\n"
        "Prefer concise, decision-ready language.\n"
    )
    if mode == "Summary":
        task = "Task: Produce a brief, bulleted summary with citations.\n"
    elif mode == "Compare":
        task = "Task: Compare key options/entities in a compact table or bullets; every row backed by citations.\n"
    elif mode == "Decision Matrix":
        task = ("Task: Build a criteria × options decision matrix from the passages only. "
                "Propose weights only if weights are explicitly present; otherwise leave weights blank. Include citations.\n")
    elif mode == "Fact Extract":
        task = ("Task: Extract key fields as JSON (also print a short human summary). "
                "Only include fields that appear in the passages. Include citations in the human summary.\n")
    elif mode == "Action Items":
        task = "Task: List concrete next actions with owner and date if present; cite each bullet.\n"
    else:
        task = "Task: Answer briefly and precisely with citations.\n"
    return baseline + style + task

def build_user_prompt(question: str, passages: List[Dict[str, Any]]) -> Tuple[str, List[str]]:
    ctx_blocks, chips = [], []
    for i, h in enumerate(passages):
        t = (h.get("text") or "").strip()
        if not t:
            continue
        ctx_blocks.append(f"[{i+1}] Source: {h.get('source')} (chunk {h.get('chunk_index')})\n{t}")
        if h.get("source"):
            chips.append(h["source"])
    context_text = "\n\n".join(ctx_blocks) if ctx_blocks else "No relevant context retrieved."
    user_prompt = (
        f"Question:\n{question}\n\n"
        f"Passages (numbered):\n{context_text}\n\n"
        "Write the result according to the Task and Rules."
    )
    return user_prompt, chips


# ---------------------- UI ----------------------
st.set_page_config(page_title="RAG Chat — Claude + Weaviate + Azure", page_icon="💬", layout="wide")

st.markdown("""
<style>
.chat-bubble-user { background:#e7f0ff; padding:10px 12px; border-radius:14px; white-space:pre-wrap; }
.chat-bubble-assistant { background:#f7f7f8; padding:10px 12px; border-radius:14px; white-space:pre-wrap; }
.small-dim { font-size:12px; opacity:.7; }
.source-chip { display:inline-block; font-size:12px; padding:4px 8px; margin:4px 6px 0 0; background:#eef1f4; border-radius:999px; }
</style>
""", unsafe_allow_html=True)

st.title("💬 RAG Chat — Claude + Weaviate + Azure Blob")

with st.sidebar:
    st.header("⚙️ Settings")
    az   = st.secrets.get("azure", {})
    weav = st.secrets.get("weaviate", {})
    anth = st.secrets.get("anthropic", {})

    # Azure
    connection_string = az.get("connection_string", "")
    container         = az.get("container", "")
    prefix            = az.get("prefix", "KB/")

    # Weaviate
    w_url     = weav.get("url", "")
    w_api_key = weav.get("api_key", None)
    w_tenancy = weav.get("tenancy", None)  # optional (v4 MT)

    # Claude
    anthropic_key = anth.get("api_key", "")

    # Retrieval knobs (non-LLM)
    cand_k  = st.slider("Candidate pool (hybrid)", 20, 120, CANDIDATES_K, step=10)
    final_k = st.slider("Final passages to LLM", 4, 20, FINAL_K_DEFAULT, step=1)
    use_reranker = st.checkbox("Use cross-encoder reranker (better precision, slower)", value=True)
    per_source_cap = st.slider("Max passages per source", 1, 6, 3, 1)
    neighbors = st.slider("Neighbor chunks to include per hit (each side)", 0, 2, 1, 1)
    hide_csv_bias = st.checkbox("Downweight CSV-like passages", value=True)

    st.divider()
    st.caption("🎛️ Answer style")
    mode = st.selectbox(
        "Mode",
        ["Direct Answer", "Summary", "Compare", "Decision Matrix", "Fact Extract", "Action Items"],
        index=0,
    )
    strictness = st.slider("Grounding strictness", 0.0, 1.0, 0.7, 0.1)

    st.divider()
    st.caption("🧠 Index KB → Weaviate")
    do_ingest = st.button("Sync / Rebuild Index")

    st.divider()
    st.caption("Health checks")
    show_health = st.checkbox("Show connection details")

    st.divider()
    st.caption("🔎 Connectivity diagnostics")
    if st.button("Run diagnostics"):
        diag = run_weaviate_diagnostics(w_url, w_api_key)
        st.json(diag)

if show_health:
    st.write("**Azure**", {"container": container, "prefix": prefix, "conn_str_set": bool(connection_string)})
    st.write("**Weaviate**", {
        "url": w_url,
        "api_key": bool(w_api_key),
        "client_version": getattr(weaviate, "__version__", "unknown"),
        "v4_mode": V4,
    })
    st.write("**Anthropic**", {"api_key": bool(anthropic_key)})

# Guards
if not w_url:
    st.error("Please set [weaviate].url in .streamlit/secrets.toml (cluster URL, no /v1)")
    st.stop()
if not anthropic_key:
    st.warning("Missing [anthropic].api_key — add your ANTHROPIC key to secrets to chat.")

# Connect
w_client = make_weaviate_client(w_url, w_api_key)
st.success("✅ Connected to Weaviate")

# Ingest
if do_ingest:
    if not (connection_string and container):
        st.error("Azure [connection_string] and [container] are required in secrets.")
    else:
        status_box = st.empty()
        def _status(msg): status_box.write(msg)
        with st.spinner("Syncing KB from Azure and (re)indexing…"):
            total = ingest_from_azure_to_weaviate(connection_string, container, prefix, w_client, tenancy=w_tenancy, status_cb=_status)
        st.success(f"Done. Upserted {total} chunks.")


# ====================== DI AGENT: Skills (deterministic, KB-only) ======================
NUM_PAT = re.compile(r"[-+]?\d{1,3}(?:,\d{3})*(?:\.\d+)?")
DATE_PAT = re.compile(r"\b(20\d{2}[-/]\d{1,2}[-/]\d{1,2}|\d{1,2}[-/]\d{1,2}[-/](20\d{2}))\b")

def _norm_num(s: str) -> Optional[float]:
    try:
        return float(s.replace(",", ""))
    except Exception:
        return None

def _collect_numbers(texts: List[str]) -> List[float]:
    vals = []
    for t in texts:
        for m in NUM_PAT.findall(t or ""):
            v = _norm_num(m)
            if v is not None:
                vals.append(v)
    return vals

def _collect_kv_candidates(texts: List[str]) -> Dict[str, Dict[str, str]]:
    """
    Extract naive key: value pairs (e.g., 'Cost: $1200', 'Latency = 50ms') grouped by 'option'
    Heuristic: lines that start with a token capitalized may be option; else try to detect 'Option A' etc.
    """
    options: Dict[str, Dict[str, str]] = {}
    cur = None
    for t in texts:
        for line in (t or "").splitlines():
            line = line.strip()
            if not line:
                continue
            if re.match(r"^(Option|Plan|Model|Vendor|Tool|Service)\s+[\w\-\.]+", line, re.I):
                cur = line.split()[0] + " " + " ".join(line.split()[1:2])
                options.setdefault(cur, {})
                continue
            m = re.match(r"^([\w\s\-/]+)[:=]\s*(.+)$", line)
            if m:
                k = m.group(1).strip()
                v = m.group(2).strip()
                if cur is None:
                    cur = "Item"
                    options.setdefault(cur, {})
                options[cur][k] = v
    return options

def _collect_dates(texts: List[str]) -> List[str]:
    ds = []
    for t in texts:
        ds.extend(DATE_PAT.findall(t or ""))
    # flatten tuples, normalize
    flat = []
    for d in ds:
        if isinstance(d, tuple):
            d = d[0] or d[1]
        flat.append(str(d))
    return flat

def skill_summary_bullets(passages: List[Dict[str, Any]]) -> str:
    bullets = []
    for i, p in enumerate(passages, 1):
        txt = (p.get("text") or "").strip()
        if not txt:
            continue
        # take first 1-2 sentences
        chunk = re.split(r"(?<=[.!?])\s+", txt)[:2]
        if chunk:
            bullets.append(f"• {chunk[0].strip()} [{i}]")
    return "\n".join(bullets[:12]) if bullets else ""

def skill_compare_table(passages: List[Dict[str, Any]]) -> Optional[str]:
    texts = [p.get("text") or "" for p in passages]
    kv = _collect_kv_candidates(texts)
    if not kv:
        return None
    # Collect columns
    all_keys = set()
    for d in kv.values():
        all_keys.update(d.keys())
    cols = ["Option"] + sorted(all_keys)
    # Build markdown table
    lines = [" | ".join(cols), " | ".join(["---"] * len(cols))]
    for opt, vals in kv.items():
        row = [opt] + [vals.get(k, "") for k in sorted(all_keys)]
        lines.append(" | ".join(row))
    return "\n".join(lines)

def skill_decision_matrix(passages: List[Dict[str, Any]]) -> Optional[str]:
    texts = [p.get("text") or "" for p in passages]
    kv = _collect_kv_candidates(texts)
    if not kv:
        return None
    crit = sorted({k for d in kv.values() for k in d.keys()})
    opts = sorted(kv.keys())
    if not crit or not opts:
        return None
    # Simple unweighted presence matrix (✓ if value present)
    lines = ["Criteria \\ Option | " + " | ".join(opts),
             " | ".join(["---"] * (len(opts) + 1))]
    for c in crit:
        row = [c]
        for o in opts:
            row.append("✓" if kv[o].get(c) else "")
        lines.append(" | ".join(row))
    return "\n".join(lines)

def skill_fact_extract_json(passages: List[Dict[str, Any]]) -> Optional[str]:
    texts = [p.get("text") or "" for p in passages]
    kv = _collect_kv_candidates(texts)
    if not kv:
        return None
    return json.dumps(kv, indent=2, ensure_ascii=False)

def _au_date(s: str) -> Optional[str]:
    # Normalize to ISO if possible; be liberal
    try:
        if "-" in s or "/" in s:
            for fmt in ("%Y-%m-%d", "%d-%m-%Y", "%Y/%m/%d", "%d/%m/%Y", "%m/%d/%Y", "%d/%m/%y"):
                try:
                    dtobj = dt.datetime.strptime(s.replace("/", "-"), fmt)
                    return dtobj.date().isoformat()
                except Exception:
                    pass
    except Exception:
        pass
    return None

def skill_action_items(passages: List[Dict[str, Any]]) -> Optional[str]:
    out = []
    for i, p in enumerate(passages, 1):
        txt = (p.get("text") or "")
        for line in txt.splitlines():
            if re.search(r"\b(action|todo|task|owner|due|deadline|assign)\b", line, re.I):
                dts = _collect_dates([line])
                iso = _au_date(dts[0]) if dts else ""
                line = re.sub(r"\s+\[\d+\]$", "", line.strip())
                out.append(f"• {line} {('— Due: ' + iso) if iso else ''} [{i}]")
    return "\n".join(out[:20]) if out else None

def skill_calc(passages: List[Dict[str, Any]]) -> Optional[str]:
    texts = [p.get("text") or "" for p in passages]
    nums = _collect_numbers(texts)
    if not nums:
        return None
    try:
        total = sum(nums)
        avg = statistics.mean(nums)
        mx = max(nums); mn = min(nums)
        return f"Computed from context numbers: count={len(nums)}, total={total:.4g}, avg={avg:.4g}, min={mn:.4g}, max={mx:.4g}"
    except Exception:
        return None


# ====================== DI AGENT: Orchestration ======================
def agent_plan(mode: str) -> List[str]:
    """Return ordered skill names to run before/alongside LLM based on mode."""
    if mode == "Summary":
        return ["summary"]
    if mode == "Compare":
        return ["compare", "calc"]
    if mode == "Decision Matrix":
        return ["decision_matrix", "calc"]
    if mode == "Fact Extract":
        return ["fact_json"]
    if mode == "Action Items":
        return ["actions"]
    return ["calc"]  # Direct Answer -> calc (optional)

def agent_act(skills: List[str], passages: List[Dict[str, Any]]) -> Dict[str, str]:
    """Run deterministic skills and return sidecar artifacts (added to final answer)."""
    artifacts: Dict[str, str] = {}
    if "summary" in skills:
        s = skill_summary_bullets(passages)
        if s: artifacts["Summary"] = s
    if "compare" in skills:
        t = skill_compare_table(passages)
        if t: artifacts["Comparison"] = t
    if "decision_matrix" in skills:
        m = skill_decision_matrix(passages)
        if m: artifacts["Decision Matrix"] = m
    if "fact_json" in skills:
        j = skill_fact_extract_json(passages)
        if j: artifacts["Facts JSON"] = j
    if "actions" in skills:
        a = skill_action_items(passages)
        if a: artifacts["Action Items"] = a
    if "calc" in skills:
        c = skill_calc(passages)
        if c: artifacts["Calculations"] = c
    return artifacts


# ---------------------- CHAT (Decision-Intelligence Agent) ----------------------
if "history" not in st.session_state:
    st.session_state.history = []
if "last_answer_payload" not in st.session_state:
    st.session_state.last_answer_payload = None  # for export

def add_history(role: str, content: str, sources: Optional[List[str]] = None):
    st.session_state.history.append({"role": role, "content": content, "sources": sources or []})

def render_msg(role: str, content: str, sources: Optional[List[str]] = None):
    css = "chat-bubble-user" if role == "user" else "chat-bubble-assistant"
    with st.chat_message(role):
        st.markdown(f"<div class='{css}'>{esc(content)}</div>", unsafe_allow_html=True)
        if role == "assistant" and sources:
            chips = "".join([f"<span class='source-chip'>{esc(s)}</span>" for s in sources[:12]])
            st.markdown(f"<div class='small-dim'>Sources: {chips}</div>", unsafe_allow_html=True)

def retrieval_quality(hits: List[Dict[str, Any]]) -> Tuple[float, int, int]:
    """Return (quality_score, total_chars, unique_sources)."""
    if not hits:
        return (0.0, 0, 0)
    total_chars = sum(len((h.get("text") or "")) for h in hits)
    uniq_sources = len({h.get("source") for h in hits if h.get("source")})
    score = min(1.0, (total_chars / 4000.0) * 0.6 + (uniq_sources / 6.0) * 0.4)
    return (score, total_chars, uniq_sources)

def gate_answerability(hits: List[Dict[str, Any]], min_chars: int = 800, min_sources: int = 2) -> bool:
    _, total_chars, uniq_sources = retrieval_quality(hits)
    return (total_chars >= min_chars) and (uniq_sources >= min_sources)

def format_context_for_prompt(hits: List[Dict[str, Any]], max_chars: int = 2000) -> Tuple[str, List[str]]:
    ctx_blocks, chips = [], []
    for i, h in enumerate(hits):
        t = (h.get("text") or "").strip()
        if not t:
            continue
        t = _truncate(t, max_chars)
        src = h.get("source")
        ctx_blocks.append(f"[{i+1}] Source: {src} (chunk {h.get('chunk_index')})\n{t}")
        if src:
            chips.append(src)
    return ("\n\n".join(ctx_blocks) if ctx_blocks else "No relevant context retrieved."), chips

def enforce_citations(text: str) -> bool:
    return has_citation(text)

def build_followups(question: str, mode: str, quality: float) -> List[str]:
    sugs = []
    if quality < 0.5:
        sugs.append("Sync the Knowledge Base or broaden the prefix to include the missing files.")
        sugs.append("Ask a narrower, factual question that appears verbatim in your docs.")
    if mode in {"Decision Matrix", "Compare"}:
        sugs.append("Provide explicit criteria and weight hints present in your KB (dates, costs, SLAs).")
    if mode == "Fact Extract":
        sugs.append("Name the exact fields you want (as they appear in your docs).")
    return sugs[:3]

st.subheader("Chat")
colL, colR = st.columns([4, 1])
with colL:
    user_q = st.chat_input("Ask about your Knowledge Base…")
with colR:
    export_md = st.button("⬇️ Export MD")
    export_json = st.button("⬇️ Export JSON")

if export_md and st.session_state.last_answer_payload:
    payload = st.session_state.last_answer_payload
    md = f"# Answer ({payload['mode']})\n\n{payload['final_text']}\n\n## Sources\n" + \
         "\n".join([f"- {s}" for s in payload.get("sources", [])])
    if payload.get("artifacts"):
        md += "\n\n## Agent Artifacts\n"
        for k, v in payload["artifacts"].items():
            md += f"\n### {k}\n\n{v}\n"
    fn = Path(tempfile.gettempdir())/"di_answer.md"
    fn.write_text(md, encoding="utf-8")
    st.success(f"Saved: {fn}")

if export_json and st.session_state.last_answer_payload:
    fn = Path(tempfile.gettempdir())/"di_answer.json"
    Path(fn).write_text(json.dumps(st.session_state.last_answer_payload, indent=2, ensure_ascii=False), encoding="utf-8")
    st.success(f"Saved: {fn}")

if user_q:
    add_history("user", user_q)
    render_msg("user", user_q)

    # Ensure collection exists and is reachable
    try:
        col = w_client.collections.get(CLASS_NAME, tenant=w_tenancy) if (V4 and w_tenancy) else w_client.collections.get(CLASS_NAME)
    except Exception as e:
        final_text = f"I can’t reach the '{CLASS_NAME}' collection. Try 'Sync / Rebuild Index' first. ({e})"
        add_history("assistant", final_text, [])
        render_msg("assistant", final_text, [])
    else:
        # 1) Plan (agent)
        plan = agent_plan(mode)

        # 2) Retrieve (strict non-LLM)
        q_vec = embed_texts([user_q])[0]
        hits = retrieve_hybrid_then_rerank(
            col,
            query_text=user_q,
            query_vec=q_vec,
            k_candidates=cand_k,
            k_final=final_k,
            use_reranker=use_reranker,
            per_source_cap=per_source_cap,
            neighbors_window=neighbors,
            downweight_csv=hide_csv_bias,
        )

        # 3) Validate answerability
        qual, total_chars, uniq_src = retrieval_quality(hits)
        enough = gate_answerability(hits)

        system_prompt = build_system_prompt(strictness, mode)
        context_text, source_chips = format_context_for_prompt(hits)

        if not enough:
            followups = build_followups(user_q, mode, qual)
            final_text = (
                "I don’t have enough information in the retrieved context to answer that.\n\n"
                f"- Retrieved characters: {total_chars}\n"
                f"- Distinct sources: {uniq_src}\n"
                + ("\n".join(f"- {s}" for s in followups) if followups else "")
            ).strip()
            add_history("assistant", final_text, source_chips)
            render_msg("assistant", final_text, source_chips)
            st.session_state.last_answer_payload = {
                "mode": mode, "question": user_q, "final_text": final_text,
                "sources": source_chips, "artifacts": {}, "answerable": False
            }
        else:
            # 4) Act (run deterministic skills pre-LLM)
            artifacts = agent_act(plan, hits)

            # 5) Compose LLM prompt and stream KB-grounded answer
            user_prompt, chips = build_user_prompt(user_q, hits)
            with st.chat_message("assistant"):
                placeholder = st.empty()
                try:
                    with call_claude(anthropic_key, system_prompt, user_prompt, stream=True) as stream_resp:
                        text_accum = ""
                        for event in stream_resp:
                            if event.type == "content_block_delta":
                                delta_text = getattr(getattr(event, "delta", None), "text", None)
                                if delta_text:
                                    text_accum += delta_text
                                    placeholder.markdown(
                                        f"<div class='chat-bubble-assistant'>{esc(text_accum)}</div>",
                                        unsafe_allow_html=True
                                    )
                        final_text = (text_accum or "").strip() or "I couldn't generate a response."
                        if not enforce_citations(final_text):
                            final_text = (
                                "I don’t have enough information in the context to answer with citations. "
                                "Please sync the KB or ask a narrower question."
                            )
                        # If we produced artifacts, append them neatly
                        if artifacts:
                            final_text += "\n\n---\n**Agent Artifacts (KB-only):**\n"
                            for k, v in artifacts.items():
                                final_text += f"\n**{k}**\n\n{v}\n"
                        placeholder.markdown(
                            f"<div class='chat-bubble-assistant'>{esc(final_text)}</div>",
                            unsafe_allow_html=True
                        )
                except Exception as e:
                    final_text = f"Sorry, streaming failed: {e}"
                    placeholder.markdown(
                        f"<div class='chat-bubble-assistant'>{esc(final_text)}</div>",
                        unsafe_allow_html=True
                    )

            add_history("assistant", final_text, source_chips)
            render_msg("assistant", final_text, source_chips)
            st.session_state.last_answer_payload = {
                "mode": mode, "question": user_q, "final_text": final_text,
                "sources": source_chips, "artifacts": artifacts, "answerable": True
            }

# History (last 8 turns)
if st.session_state.history:
    st.markdown("### Recent conversation")
    for msg in st.session_state.history[-8:]:
        render_msg(msg["role"], msg["content"], msg.get("sources"))

# Clean up v4 client
if V4:
    try:
        w_client.close()
    except Exception:
        pass
